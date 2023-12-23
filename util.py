import string
import colorsys
import numpy as np


def get_slide_notes(slide):
  result = ""
  num_breaks = 0

  if slide.notes_slide:
    for paragraph in slide.notes_slide.notes_text_frame.paragraphs:
      for run in paragraph.runs:
        if "[Break]" in run.text:
          num_breaks += 1
          continue
        elif '[' in run.text and ']' in run.text:
          continue
        result += run.text + "\n"

  return result, num_breaks


def print_all_text(prs):
  # text_runs will be populated with a list of strings,
  # one for each text run in presentation
  text_runs = []

  for slide in prs.slides:
    for shape in slide.shapes:
      if not shape.has_text_frame:
        continue
      for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
          text_runs.append(run.text)

  print(text_runs)


def is_backup_slide(slide):
  if slide.shapes.title:
    title = slide.shapes.title.text
    title = title.lower()
    if "backup" in title:
      return True
  return False


def within_bounds(prev_shape_pos, curr_shape_pos, shape_pos_threshold, slide_width, slide_height):
  prev_shape_pos_x = prev_shape_pos[0]
  prev_shape_pos_y = prev_shape_pos[1]
  curr_shape_pos_x = curr_shape_pos[0]
  curr_shape_pos_y = curr_shape_pos[1]

  step_x = slide_width * shape_pos_threshold
  step_y = slide_height * shape_pos_threshold

  shape_min_x = prev_shape_pos_x - step_x
  shape_max_x = prev_shape_pos_x + step_x
  shape_min_y = prev_shape_pos_y - step_y
  shape_max_y = prev_shape_pos_y + step_y

  return (curr_shape_pos_x > shape_min_x and curr_shape_pos_x < shape_max_x and
          curr_shape_pos_y > shape_min_y and curr_shape_pos_y < shape_max_y)


# Solution from https://groups.google.com/g/python-pptx/c/iTaK8if8Dck
def get_color_scheme(prs):
  from pptx.enum.dml import MSO_FILL, MSO_COLOR_TYPE
  from pptx.opc.constants import RELATIONSHIP_TYPE as RT
  from pptx.oxml import parse_xml
  from pptx.oxml.ns import _nsmap

  for slide in prs.slides:
    presentation_part = prs.part
    theme_part = presentation_part.part_related_by(RT.THEME)
    theme_element = parse_xml(theme_part.blob)
    xpath = 'a:themeElements/a:clrScheme'
    color_scheme = theme_element.xpath(xpath)[0]
  return(color_scheme)


def get_scheme_color_rgb(color_scheme, theme_color, brightness, debug=False):
  from pptx.enum.dml import MSO_THEME_COLOR
  theme_color_wrap_around = int(MSO_THEME_COLOR.FOLLOWED_HYPERLINK)
  if theme_color >= MSO_THEME_COLOR.DARK_1:
    entry = color_scheme[theme_color - theme_color_wrap_around - 1][0]
  elif theme_color == 0:
    entry = color_scheme[theme_color][0]
  else:
    entry = color_scheme[theme_color - 1][0]
  if entry.tag[-6:] == 'sysClr':
    color_rgb = entry.get('lastClr')
  elif entry.tag[-7:] == 'srgbClr':
    color_rgb = entry.get('val')

  if brightness:
    color_code = get_hex_code(color_rgb)
    srgb = np.array(color_code)
    srgb = srgb / 255
    h, luminance, s = colorsys.rgb_to_hls(*srgb)
    if brightness > 0:
      luminance = luminance * (1 - brightness) + (brightness)
    else:
      luminance = luminance * (1 + brightness)
    srgb = np.array(colorsys.hls_to_rgb(h, luminance, s))
    srgb = (srgb * 255).round(0).astype(int)
    rgb2hex = lambda r,g,b: '%02X%02X%02X' %(r,g,b)
    return rgb2hex(*srgb)

  return color_rgb


def get_hex_code(color):
  return [int(color[i:i+2], 16) for i in (0, 2, 4)]


def calculate_luminace(color_code):
  index = float(color_code) / 255

  if index < 0.03928:
    return index / 12.92
  else:
    return ( ( index + 0.055 ) / 1.055 ) ** 2.4


# Solution from https://github.com/Peter-Slump/python-contrast-ratio
def calculate_relative_luminance(rgb):
  return 0.2126 * calculate_luminace(rgb[0]) +
         0.7152 * calculate_luminace(rgb[1]) +
         0.0722 * calculate_luminace(rgb[2])


# Solution from https://github.com/Peter-Slump/python-contrast-ratio
def calculate_contrast_ratio(colorA, colorB):
  colorA = get_hex_code(colorA)
  colorB = get_hex_code(colorB)

  light = colorA if sum(colorA) > sum(colorB) else colorB
  dark = colorA if sum(colorA) < sum(colorB) else colorB

  contrast_ratio = ( calculate_relative_luminance(light) + 0.05 ) / ( calculate_relative_luminance(dark) + 0.05 )
  return contrast_ratio


def load_words(part_type):
  result = []
  with open("./data/%s.txt" % part_type) as f:
    for word in f:
      result.append(word.strip())
  return result


class WordSet:
  def __init__(self, wn, verbs, prepositions, articles):
    self.wn = wn
    self.verbs = verbs
    self.prepositions = prepositions
    self.articles = articles


def initialize_word_set():
  from nltk.corpus import wordnet as wn
  verbs = load_words("verbs")
  prepositions = load_words("prepositions")
  articles = ['the', 'a', 'an']
  return WordSet(wn, verbs, prepositions, articles)


def convert_string_into_word_tokens(wordtext):
  wordtext = wordtext.strip()
  wordtext = wordtext[0].lower() + wordtext[1:]
  wordtext = wordtext.replace('i.e.', '')
  wordtext = wordtext.replace('e.g.', '')
  for punctuation in string.punctuation:
    wordtext = wordtext.replace(punctuation, '')
  return wordtext.strip().split(' ')


def identify_parts_of_speech(ws, word_list):
  wn = ws.wn
  verbs = ws.verbs
  prepositions = ws.prepositions
  articles = ws.articles

  pos_all = dict()
  for w in word_list:
    pos_l = set()
    for tmp in wn.synsets(w):
      if tmp.name().split('.')[0] == w:
        pos_l.add(tmp.pos())
    if w in prepositions:
      pos_l.add('p')
    elif w in articles:
      pos_l.clear()
      pos_l.add('at')
    elif w in verbs:
      pos_l.clear()
      pos_l.add('v')
      if w.endswith("ing"):
        pos_l.add('n')
    elif w.endswith("ing"):
      pos_l.add('a')
      pos_l.add('n')
      pos_l.add('v')
    elif w.endswith("ed"):
      pos_l.clear()
      pos_l.add('a')
    elif not pos_l:
      pos_l.add('n')
    pos_all[w] = pos_l
  return pos_all


def is_full_sentence(classified_words):
  # Find patterns where: 'n','v','n'
  word_count = 0
  sentence_comp_count = 0

  for word, part in classified_words.items():
    if sentence_comp_count == 0:
      if 'v' in part:
        return False
      elif 'n' in part:
        sentence_comp_count = 1
    else:
      if 'v' in part and 'n' not in part and sentence_comp_count%2==1:
        sentence_comp_count += 1
      elif 'n' in part and 'v' not in part and sentence_comp_count%2==0:
        sentence_comp_count += 1
    word_count += 1
  return True if sentence_comp_count >= 3 else False


def display_comments_on_webpage(time_estimate, slide_feedback, slide_times,
                                cumul_slide_times, general_feedback, pass_all_checks):
  from pretty_html_table import build_table
  import webbrowser
  import pandas as pd

  df = pd.DataFrame(slide_feedback, columns=["Feedback"])
  df.index += 1
  df = df.rename_axis("Slide #").reset_index()
  df["Time at Slide Start"] = cumul_slide_times
  df["Time Spent on Slide"] = slide_times

  html_table_blue_light = build_table(df, 'blue_dark', index=False, escape=False)
  with open('output.html', 'w') as f:
    f.write("<h3>General Feedback:</h3>")
    if time_estimate:
      f.write("Estimate total time for presentation: %s" % time_estimate)

    if not general_feedback and pass_all_checks:
      f.write("<p>Presentation passed all checks!</p>")
    else:
      f.write("<p>" + general_feedback + "</p>")

    f.write(html_table_blue_light)

  webbrowser.open_new_tab("output.html")

