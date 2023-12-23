from util import *
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE, MSO_AUTO_SHAPE_TYPE
from pptx.enum.dml import MSO_FILL, MSO_COLOR_TYPE, MSO_THEME_COLOR
from pptx.util import Pt
import time


def must_end_with_summary_slide(prs):
  summary_at_end = False
  for slide in prs.slides:
    if slide.shapes.title:
      title = slide.shapes.title.text.lower()
      if "summary" in title:
        summary_at_end = True
      elif "backup" in title:
        return summary_at_end
      elif summary_at_end:
        return False
  return summary_at_end


def should_have_slide_numbers(prs, slide_feedback):
  has_slide_numbers = False
  shape_left = 0
  shape_top = 0

  slide_num = 1
  slide_height = prs.slide_height

  if len(prs.slides) < 2:
    return True

  for slide in prs.slides:

    # Skip title slide
    if slide_num == 1:
      slide_num = 2
      continue

    slide_has_slide_number = False

    if is_backup_slide(slide):
      return has_slide_numbers

    for shape in slide.shapes:
      if not shape.has_text_frame:
        continue
      num_paragraphs = len(shape.text_frame.paragraphs)

      # Mark as candidate for slide number
      if num_paragraphs == 1:
        shape_text = shape.text.strip()
        if (shape_text.isdigit() and int(shape_text) == slide_num) or shape_text=="‹#›":
          if (shape_top > slide_height * 0.9 or shape_top < slide_height * 0.1):
            slide_has_slide_number = True
            has_slide_numbers = True
            if (shape_left == 0 and shape_top == 0):
              shape_left = shape.left
              shape_top = shape.top
            elif (shape_left != shape.left or shape_top != shape.top):
              slide_feedback[slide_num - 1] += "Slide number is misplaced in a different location.\n"

    if has_slide_numbers and not slide_has_slide_number:
      slide_feedback[slide_num - 1] += "Slide number is missing on this slide.\n"

    slide_num += 1

  return has_slide_numbers


def has_smooth_slide_transitions(prs, slide_feedback):
  shape_pos_threshold = 0.1

  has_smooth_transitions = True
  shape_left = 0
  shape_top = 0

  shapes_prev = {}
  shapes_attr_prev = {}

  slide_num = 1

  if len(prs.slides) < 2:
    return True

  for slide in prs.slides:

    if is_backup_slide(slide):
      return has_smooth_transitions

    shapes_curr = {}
    shapes_attr_curr = {}

    for shape in slide.shapes:
      shape_left = shape.left
      shape_top = shape.top

      shape_index = 0

      while True:
        shape_attrs = [str(shape.shape_type), str(shape.width), str(shape.height), str(shape_index)]
        shape_hash = hash(','.join(shape_attrs))

        if shape_hash not in shapes_curr:
          break
        shape_index += 1

      shapes_curr[shape_hash] = (shape.left, shape.top)
      if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
        shapes_attr_curr[shape_hash] = [str(shape.auto_shape_type)]
      else:
        shapes_attr_curr[shape_hash] = [str(shape.shape_type)]

      if shape.has_text_frame and shape.text:
        shapes_attr_curr[shape_hash].append(shape.text)

    if shapes_prev and shapes_curr:
      for shape_hash, curr_shape_pos in shapes_curr.items():
        if shape_hash in shapes_prev:
          prev_shape_pos = shapes_prev[shape_hash]
          if (prev_shape_pos == curr_shape_pos and
              len(shapes_attr_curr[shape_hash]) == len(shapes_attr_prev[shape_hash])):
            del shapes_prev[shape_hash]

      for shape_hash, prev_shape_pos in shapes_prev.items():
        if shape_hash in shapes_curr:
          curr_shape_pos = shapes_curr[shape_hash]
          shape_attr = shapes_attr_curr[shape_hash]

          # Only match those with either the same shape or same shape and text
          if (len(shape_attr) == len(shapes_attr_prev[shape_hash]) and
              (prev_shape_pos != curr_shape_pos and
               within_bounds(prev_shape_pos, curr_shape_pos,
                             shape_pos_threshold, prs.slide_width, prs.slide_height)) and
              shapes_attr_prev[shape_hash] == shape_attr):
            has_smooth_transitions = False
            if len(shape_attr) > 1:
              if (shape_attr[1] == "‹#›"):
                slide_feedback_comment = "Slide transition for the slide number is not smooth.\n"
              else:
                slide_feedback_comment = "Slide transition for %s is not smooth. This shape object holds the following text: '%s'\n" % (shape_attr[0], shape_attr[1])
            else:
              slide_feedback_comment = "Slide transition for %s is not smooth.\n" % shape_attr[0]
              slide_feedback[slide_num - 1] += slide_feedback_comment

    slide_num += 1
    shapes_prev = shapes_curr
    shapes_attr_prev = shapes_attr_curr

  return has_smooth_transitions


def should_have_high_contrast_fonts_colours(prs, slide_feedback):
# Only checks colours of shapes, textboxes, lines, but not pictures and graphs
  shape_min_color_contrast_ratio = 1.5
  font_min_color_contrast_ratio = 3
  min_size_font = 18
  min_line_width = 2

  slide_num = 1
  result = True

  color_scheme = get_color_scheme(prs)

  for slide in prs.slides:

    if is_backup_slide(slide):
      return result

    shapes_curr = {}
    shapes_attr_curr = {}

    if slide.background.fill.type != MSO_FILL.SOLID:
      slide_background_color = "FFFFFF"
    elif slide.background.fill.fore_color == MSO_COLOR_TYPE.RGB:
      slide_background_color = slide.background.fill.fore_color.rgb.__str__()

    for shape in slide.shapes:
      shape_type = shape.shape_type
      shape_feedback_comment_temp = ""

      if (shape_type == MSO_SHAPE_TYPE.PICTURE or
          shape_type == MSO_SHAPE_TYPE.CHART or
          shape_type == MSO_SHAPE_TYPE.TABLE):
        continue

      if shape_type == MSO_SHAPE_TYPE.LINE:
        shape = shape.line
        line_width = shape.width.pt
        if line_width < min_line_width:
          slide_feedback_comment = "Line width for %s is too small to be seen at %s pts.\n" % (str(shape_type), str(shape.width.pt))
          slide_feedback[slide_num - 1] += slide_feedback_comment
          result = False

      fill_format = shape.fill

      font_check_against_color = slide_background_color
      at_least_one_font_visible = False # Some fonts may be intentionally greyed out

      # Only check fills of shapes that have a solid fill
      if fill_format.type == MSO_FILL.SOLID:
        color_format = fill_format.fore_color
        if color_format.type == MSO_COLOR_TYPE.RGB:
          color_rgb = color_format.rgb.__str__()
        else:
          color_rgb = get_scheme_color_rgb(color_scheme,
                                           shape.fill.fore_color.theme_color,
                                           shape.fill.fore_color.brightness, True)

        is_rectangle = False
        contrast_ratio = calculate_contrast_ratio(slide_background_color, color_rgb)
        if (contrast_ratio < shape_min_color_contrast_ratio and not contrast_ratio==1):
          if shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            shape_descriptor = shape.auto_shape_type
            if (shape_descriptor == MSO_SHAPE.RECTANGLE or
                shape_descriptor == MSO_SHAPE.ROUNDED_RECTANGLE):
              is_rectangle = True
          else:
            shape_descriptor = shape_type

          # Ignore rectangles as they are often used to cover components
          if not is_rectangle:
            slide_feedback_comment = "Colour contrast for %s is not sufficient from the slide background colour.\n" % str(shape_descriptor)
            slide_feedback[slide_num - 1] += slide_feedback_comment
            result = False

        font_check_against_color = color_rgb

      if shape_type != MSO_SHAPE_TYPE.LINE and shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
          for run in paragraph.runs:
            font = run.font

            if (font.size and
                ((font.size < Pt(min_size_font) and
                  len(run.text.split()) > 2 and
                  not run.text.startswith('*')) or
                 (font.size < Pt(min_size_font - 6)))):
              if shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                shape_descriptor = shape.auto_shape_type
              else:
                shape_descriptor = shape_type
              slide_feedback_comment = "Font size for text '%s' in shape %s is too small.\n" % (run.text, str(shape_descriptor))
              slide_feedback[slide_num - 1] += slide_feedback_comment
              result = False

            if not len(run.text):
              continue

            if font.color.type == MSO_COLOR_TYPE.RGB:
              font_color_rgb = font.color.rgb.__str__()
            else:
              # Set to default font if font color is not found
              if not font.color.type:
                font.color.theme_color = MSO_THEME_COLOR.DARK_1

              font_color_rgb = get_scheme_color_rgb(color_scheme,
                                                    font.color.theme_color,
                                                    font.color.brightness)

            contrast_ratio = calculate_contrast_ratio(font_check_against_color, font_color_rgb)
            if (contrast_ratio < font_min_color_contrast_ratio and not at_least_one_font_visible):
              if shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                shape_descriptor = shape.auto_shape_type
              else:
                shape_descriptor = shape_type
              shape_feedback_comment_temp += "Font colour contrast for text '%s' in shape %s is not sufficient from the background colour.\n" % (run.text, str(shape_descriptor))
            else:
              at_least_one_font_visible = True # Change back to true

      if not at_least_one_font_visible and shape_feedback_comment_temp:
        slide_feedback[slide_num - 1] += shape_feedback_comment_temp
        result = False

    slide_num += 1

  return result


def should_not_have_excessive_text(prs, slide_feedback):
  max_num_words_per_slide = 30

  has_excessive_text = False
  slide_num = 1

  for slide in prs.slides:
    slide_text = ""

    for shape in slide.shapes:
      if not shape.has_text_frame:
        continue
      for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
          if (len(run.text.split(' ')) > 2):
            if slide.shapes.title:
              title = slide.shapes.title.text
              if run.text.strip() == title.strip():
                continue
            slide_text += run.text.strip() + " "
    word_count = len(slide_text.split(' '))

    if word_count > max_num_words_per_slide:
      slide_feedback_comment = "Excessive amount of words on this slide.\n"
      slide_feedback[slide_num - 1] += slide_feedback_comment
      has_excessive_text = True

    slide_num += 1

  return not has_excessive_text


def does_not_have_complete_sentences(prs, slide_feedback):
  result = True
  wordset = initialize_word_set()
  slide_num = 1

  for slide in prs.slides:
    title = ""
    if slide.shapes.title:
      title = slide.shapes.title.text.lower()

    for shape in slide.shapes:
      if not shape.has_text_frame:
        continue
      for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
          shape_text = run.text.strip()

          if (shape_text and shape_text.strip() != title.strip() and
              len(shape_text.split(' ')) > 4 and
              not shape_text.endswith('?') and
              ':' not in shape_text and '-' not in shape_text):
            word_tokens = convert_string_into_word_tokens(shape_text)
            words_classified = identify_parts_of_speech(wordset, word_tokens)
            if is_full_sentence(words_classified):
              slide_feedback_comment = "Avoid full sentences: '%s'\n" % run.text
              slide_feedback[slide_num - 1] += slide_feedback_comment
              result = False

    slide_num += 1

  return result


def estimate_presentation_length(prs):
  seconds_per_word = 0.42
  seconds_per_pause = 1.5
  seconds_per_break = 3
  seconds_between_slides = 2

  total_prs_time = 0
  string_punctuation = ['.', '?', '!']
  slide_num = 1
  slides_without_notes = 0
  slide_times = []
  cumul_slide_times = []

  for slide in prs.slides:
    slide_notes, num_breaks = get_slide_notes(slide)
    if not slide_notes:
      slides_without_notes += 1
    if slides_without_notes > 2:
      return None, None, None

    if slide.shapes.title:
      title = slide.shapes.title.text.lower()
      if "backup" in title:
        break

    cumul_slide_times.append(time.strftime('%H:%M:%S', time.gmtime(total_prs_time)))
    time_per_slide = num_breaks * seconds_per_break

    for punctuation in string_punctuation:
      punc_pauses = slide_notes.count(punctuation)
      time_per_slide += punc_pauses * seconds_per_pause
      slide_notes = slide_notes.replace(punctuation, '')

    words = slide_notes.strip().split(' ')
    time_per_slide += len(words) * seconds_per_word

    slide_num += 1
    total_prs_time += time_per_slide + seconds_between_slides
    slide_times.append(time.strftime('%M:%S', time.gmtime(time_per_slide)))

  return time.strftime('%H:%M:%S', time.gmtime(total_prs_time)), slide_times, cumul_slide_times

